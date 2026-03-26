#!/usr/bin/env python3
"""
pptx2md_llm.py – 将 PPTX 转换为 LLM 友好的 Markdown 格式

用法:
    python pptx2md_llm.py input.pptx                    # 输出到 stdout
    python pptx2md_llm.py input.pptx -o output.md       # 输出到文件
    python pptx2md_llm.py input.pptx -o out.md -i imgs/  # 同时导出图片

依赖:
    pip install python-pptx
"""

import argparse
import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


# ──────────────────── 工具函数 ────────────────────

def emu_to_cm(emu_val):
    """EMU (English Metric Unit) 转厘米，保留 1 位小数"""
    if emu_val is None:
        return "?"
    return round(emu_val / 914400 * 2.54, 1)


def position_tag(shape):
    """生成简洁的位置描述 [位置: left×top, size宽×高 cm]"""
    l = emu_to_cm(shape.left)
    t = emu_to_cm(shape.top)
    w = emu_to_cm(shape.width)
    h = emu_to_cm(shape.height)
    return f"[位置: ({l}, {t}) cm, 尺寸: {w}×{h} cm]"


def is_title_shape(shape):
    """判断 shape 是否为标题占位符"""
    if not shape.is_placeholder:
        return False
    ph_type = shape.placeholder_format.type
    return ph_type in (
        PP_PLACEHOLDER.TITLE,
        PP_PLACEHOLDER.SUBTITLE,
        PP_PLACEHOLDER.CENTER_TITLE,
        PP_PLACEHOLDER.VERTICAL_TITLE,
    )


# ──────────────────── 文本提取 ────────────────────

def extract_paragraph_text(paragraph):
    """
    提取一个段落的纯文本，保留分点和编号信息。
    返回 (text, level, is_bullet)
    """
    text = ""
    for run in paragraph.runs:
        text += run.text
    # 如果 paragraph.runs 为空但有文本，也尝试获取
    if not text:
        text = paragraph.text or ""
    text = text.strip()
    level = paragraph.level if paragraph.level else 0
    # 判断是否是列表项（有缩进层级 或 有 bullet/numbering）
    is_bullet = level > 0
    # 检查是否有 bullet XML 标记
    pPr = paragraph._p.find(
        '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr'
    )
    if pPr is not None:
        buNone = pPr.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}buNone'
        )
        buChar = pPr.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}buChar'
        )
        buAutoNum = pPr.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum'
        )
        if buNone is not None:
            is_bullet = False
        elif buChar is not None or buAutoNum is not None:
            is_bullet = True

    return text, level, is_bullet


def format_text_block(shape, include_position=True):
    """将文本框的内容转为 Markdown 文本"""
    lines = []
    if include_position:
        lines.append(position_tag(shape))
    for para in shape.text_frame.paragraphs:
        text, level, is_bullet = extract_paragraph_text(para)
        if not text:
            continue
        indent = "  " * level
        if is_bullet:
            lines.append(f"{indent}* {text}")
        else:
            lines.append(f"{indent}{text}")
    return "\n".join(lines)


# ──────────────────── 表格提取 ────────────────────

def format_table(shape, include_position=True):
    """将表格 shape 转为 Markdown 表格"""
    table = shape.table
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            cell_text = cell.text.replace("\n", " ").strip()
            cells.append(cell_text)
        rows.append(cells)

    if not rows:
        return ""

    lines = []
    if include_position:
        lines.append(position_tag(shape))

    # 表头
    lines.append("| " + " | ".join(rows[0]) + " |")
    lines.append("| " + " | ".join(["---"] * len(rows[0])) + " |")
    # 数据行
    for row in rows[1:]:
        # 处理合并单元格导致的列数不一致
        while len(row) < len(rows[0]):
            row.append("")
        lines.append("| " + " | ".join(row[:len(rows[0])]) + " |")

    return "\n".join(lines)


# ──────────────────── 图片提取 ────────────────────

def format_image(shape, slide_idx, img_counter, image_dir=None, include_position=True):
    """
    处理图片 shape，返回 (markdown_text, new_counter)。
    如果提供 image_dir 则导出图片文件。
    """
    parts = []
    if include_position:
        parts.append(position_tag(shape))

    if image_dir:
        try:
            img = shape.image
            ext = img.ext or "png"
            filename = f"slide{slide_idx + 1}_img{img_counter}.{ext}"
            filepath = Path(image_dir) / filename
            filepath.parent.mkdir(parents=True, exist_ok=True)
            with open(filepath, "wb") as f:
                f.write(img.blob)
            parts.append(f"![图片]({filepath})")
        except Exception:
            parts.append("[图片: 无法导出]")
    else:
        # 无输出目录时仅做文字标记
        try:
            ext = shape.image.ext or "image"
            w = emu_to_cm(shape.width)
            h = emu_to_cm(shape.height)
            parts.append(f"[图片: {ext}格式, {w}×{h} cm]")
        except Exception:
            parts.append("[图片]")

    return "\n".join(parts), img_counter + 1


# ──────────────────── 递归解组 ────────────────────

def ungroup_shapes(shapes):
    """递归展开分组的 shapes"""
    result = []
    for shape in shapes:
        try:
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                result.extend(ungroup_shapes(shape.shapes))
            else:
                result.append(shape)
        except Exception:
            result.append(shape)
    return result


# ──────────────────── 单页处理 ────────────────────

# 短文本判定阈值：纯文本字数 <= 此值 且 无列表标记 的文本框视为"标签"
LABEL_MAX_CHARS = 15


def _is_label_shape(shape):
    """判断文本框是否只是图表里的小标签"""
    text = shape.text_frame.text.strip().replace("\n", "")
    if len(text) > LABEL_MAX_CHARS:
        return False
    # 如果含有列表项 / 多段有层级，视为正文
    for para in shape.text_frame.paragraphs:
        if para.level and para.level > 0:
            return False
    return True


def process_slide(slide, slide_idx, image_dir=None, include_position=True):
    """
    处理单张幻灯片，返回 Markdown 字符串。
    """
    # 展开并按位置排序 (从上到下，从左到右)
    shapes = ungroup_shapes(slide.shapes)
    shapes_sorted = sorted(shapes, key=lambda s: (s.top or 0, s.left or 0))

    slide_title = None
    content_blocks = []
    labels = []  # 收集短标签，最后聚合输出

    for shape in shapes_sorted:
        # ── 标题 ──
        if is_title_shape(shape):
            title_text = shape.text_frame.text.strip() if shape.has_text_frame else ""
            if title_text:
                slide_title = title_text
            continue

        # ── 表格 ──
        if shape.has_table:
            md = format_table(shape, include_position)
            if md:
                content_blocks.append(md)
            continue

        # ── 图片和形状：跳过 ──
        if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE,
                                MSO_SHAPE_TYPE.AUTO_SHAPE,
                                MSO_SHAPE_TYPE.FREEFORM,
                                MSO_SHAPE_TYPE.LINE):
            continue

        # ── 文本框 ──
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not text:
                continue
            # 短文本 → 归入标签
            if _is_label_shape(shape):
                label = text.replace("\n", " ")
                labels.append(label)
            else:
                md = format_text_block(shape, include_position)
                if md.strip():
                    content_blocks.append(md)
            continue

        # 其他类型直接跳过

    # 组装本页 Markdown
    page_num = slide_idx + 1
    header = f"## 第{page_num}页：{slide_title or '(无标题)'}"

    # 提取备注
    notes_text = ""
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            notes_text = f"\n> **备注:** {notes}"

    body = "\n\n".join(content_blocks)

    # 去重标签并聚合为一行
    if labels:
        seen = set()
        unique = []
        for lb in labels:
            if lb not in seen:
                seen.add(lb)
                unique.append(lb)
        label_line = "标签：" + " | ".join(unique)
    else:
        label_line = ""

    parts = [header]
    if body:
        parts.append(body)
    if label_line:
        parts.append(label_line)
    if notes_text:
        parts.append(notes_text)

    return "\n\n".join(parts)


# ──────────────────── 主转换函数 ────────────────────

def pptx_to_markdown(pptx_path, image_dir=None, include_position=False):
    """
    将 PPTX 文件转换为 LLM 友好的 Markdown 文本。

    Args:
        pptx_path: PPTX 文件路径
        image_dir: 图片导出目录 (None 则不导出图片文件，仅做文字标注)
        include_position: 是否在输出中包含元素位置信息

    Returns:
        Markdown 格式的字符串
    """
    prs = Presentation(pptx_path)

    # 幻灯片尺寸
    slide_w = emu_to_cm(prs.slide_width)
    slide_h = emu_to_cm(prs.slide_height)

    meta_lines = [
        f"# {Path(pptx_path).stem}",
        "",
        f"幻灯片尺寸：{slide_w} × {slide_h} cm | 共 {len(prs.slides)} 页",
        "",
        "---",
        "",
    ]

    slide_mds = []
    for idx, slide in enumerate(prs.slides):
        md = process_slide(slide, idx, image_dir, include_position)
        slide_mds.append(md)

    return "\n".join(meta_lines) + "\n\n---\n\n".join(slide_mds) + "\n"


# ──────────────────── CLI ────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="将 PPTX 转换为 LLM 友好的 Markdown 格式"
    )
    parser.add_argument("pptx", type=Path, help="输入的 PPTX 文件路径")
    parser.add_argument(
        "-o", "--output", type=Path, default=None,
        help="输出 Markdown 文件路径 (默认输出到 stdout)"
    )
    parser.add_argument(
        "-i", "--image-dir", type=Path, default=None,
        help="图片导出目录 (不指定则仅文字描述图片)"
    )
    parser.add_argument(
        "--with-position", action="store_true",
        help="输出元素位置信息（默认关闭）"
    )
    args = parser.parse_args()

    if not args.pptx.is_file():
        print(f"错误: 文件不存在 – {args.pptx}", file=sys.stderr)
        sys.exit(1)

    md = pptx_to_markdown(
        str(args.pptx),
        image_dir=str(args.image_dir) if args.image_dir else None,
        include_position=args.with_position,
    )

    if args.output:
        args.output.parent.mkdir(parents=True, exist_ok=True)
        with open(args.output, "w", encoding="utf-8") as f:
            f.write(md)
        print(f"✅ 已写入 {args.output}", file=sys.stderr)
    else:
        print(md)


if __name__ == "__main__":
    main()
